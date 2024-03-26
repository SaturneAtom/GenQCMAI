from dotenv import load_dotenv
import os
from PyPDF2 import PdfReader
import streamlit as st
from langchain.text_splitter import CharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores.faiss import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain_openai import OpenAI
from langchain_community.callbacks import get_openai_callback
import requests
import time
from pptx import Presentation
from io import BytesIO
from PIL import Image
from pptx.util import Inches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load environment variables
load_dotenv()

# Déclaration des variables globales pour stocker les résultats de génération
first_generation_result = None
second_generation_result = None
three_generation_result = None
bullet_points = None
title = None


def process_text(text):
    # Split the text into chunks using Langchain's CharacterTextSplitter
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=300,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    
    # Convert the chunks of text into embeddings to form a knowledge base
    embeddings = OpenAIEmbeddings()
    knowledgeBase = FAISS.from_texts(chunks, embeddings)
    
    return knowledgeBase

def main():
    global first_generation_result, second_generation_result, description_generation_result, bullet_points, title
    
    description_generation_result = None  # Définir la variable ici
    
    # Define queries
    # Générer image avec Dall-" prompt"
    query_description_powerpoint = "Imaginez-vous être un professeur d'université enseignant un cours basé sur le document fourni. Votre objectif est de fournir une description pour une image avec Dall-E qui pourrait aider les étudiants à mieux comprendre le contenu."
    query_list_questions = "Imaginez-vous être un professeur d'université enseignant un cours basé sur le document fourni. Votre objectif est d'aider les étudiants à mieux comprendre le contenu en formulant des questions pertinentes. Générez une liste de 7 questions qui pourraient être posées aux étudiants pour les guider dans leur apprentissage. Assurez-vous que les questions couvrent les concepts clés, les exemples spécifiques et encouragent une réflexion approfondie sur le sujet"
    query_text_trou = "Imaginez-vous être un professeur d'université enseignant un cours basé sur le document fourni. Votre objectif est d'aider les étudiants à mieux comprendre le contenu en formulant des texte à trous pertinentes."
    query_title_powerpoint = "Imaginez-vous être un professeur d'université enseignant un cours basé sur le document fourni. Votre objectif est de créer un titre de PowerPoint pour aider les étudiants à mieux comprendre le contenu. En français, le titre d'une diapositive PowerPoint est une phrase courte qui résume le contenu de la diapositive."
    query_bullet_point_powerpoint = "Imaginez-vous être un professeur d'université enseignant un cours basé sur le document fourni. Votre objectif est de créer des points à puces pour aider les étudiants à mieux comprendre le contenu. En français, les points à puces sont des listes à puces qui permettent de structurer les informations de manière claire et concise."
    
    st.title("Chat with your PDF 💬")
    
    pdf = st.file_uploader('Upload your PDF Document', type='pdf', key="pdf_uploader")
    
    if pdf is not None:
        pdf_reader = PdfReader(pdf)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        
        knowledgeBase = process_text(text)
        
        st.write("PDF uploaded successfully! 🎉")

        tab1, tab2, tab3 = st.tabs(["Créer une liste de 7 question", "Créer un exercice de texte à trous", "Owl"])

        with tab1:    
            if st.button("Générer liste de questions", key="list_questions_key"):
                if first_generation_result is None:
                    docs = knowledgeBase.similarity_search(query_list_questions)
                    llm = OpenAI()
                    chain = load_qa_chain(llm, chain_type='stuff')
                    with get_openai_callback() as cost:
                        first_generation_result = chain.run(input_documents=docs, question=query_list_questions)
                        print(cost)
                st.write(first_generation_result)

        with tab2:
            if st.button("Générer texte à trous", key="text_trou_key"):
                if second_generation_result is None:
                    docs = knowledgeBase.similarity_search(query_text_trou)
                    llm = OpenAI()
                    chain = load_qa_chain(llm, chain_type='stuff')
                    with get_openai_callback() as cost:
                        second_generation_result = chain.run(input_documents=docs, question=query_text_trou)
                        print(cost)
                st.write(second_generation_result)
        
        with tab3:
            if st.button("Générer PowerPoint", key="create_powerpoint"):
                if description_generation_result is None:
                    docs = knowledgeBase.similarity_search(query_description_powerpoint)
                    llm = OpenAI()
                    chain = load_qa_chain(llm, chain_type='stuff')
                    with get_openai_callback() as cost:
                        description_generation_result = chain.run(input_documents=docs, question=query_description_powerpoint)
                        print(cost)
                time.sleep(10)

                if title is None:
                    docs = knowledgeBase.similarity_search(query_title_powerpoint)
                    llm = OpenAI()
                    chain = load_qa_chain(llm, chain_type='stuff')
                    with get_openai_callback() as cost:
                        title = chain.run(input_documents=docs, question=query_title_powerpoint)
                        print(cost)
                st.write(title)

                if bullet_points is None:
                    docs = knowledgeBase.similarity_search(query_bullet_point_powerpoint)
                    llm = OpenAI()
                    chain = load_qa_chain(llm, chain_type='stuff')
                    with get_openai_callback() as cost:
                        bullet_points = chain.run(input_documents=docs, question=query_bullet_point_powerpoint)
                        print(cost)
                st.write(bullet_points)
                time.sleep(20)

                if description_generation_result is not None:
                    response = requests.post(
                        "https://api.openai.com/v1/images/generations",
                        headers={
                            "Authorization": f"Bearer {os.getenv('OPENAI_API_KEY')}",
                            "Content-Type": "application/json"
                        },
                        json={
                            "model": "dall-e-3",
                            "prompt": f"{description_generation_result}",
                            "n": 1,
                            "size": "1024x1024"
                        }
                    )

                    print(response.json())

                    if response.status_code == 200:
                        image_url = response.json()["data"][0]["url"]
                        st.image(image_url)
                        # Créer la présentation PowerPoint
                        prs = Presentation()

                        # Diapositive avec titre à gauche et image à droite
                        slide_layout = prs.slide_layouts[5]  # Utiliser un layout adapté
                        slide = prs.slides.add_slide(slide_layout)

                        title_shape = slide.placeholders[0]
                        title_shape.text = title
                        title_shape.text_frame.paragraphs[0].font.size = Pt(16)
                        response_img = requests.get(image_url)
                        if response_img.status_code == 200:
                            image_data = BytesIO(response_img.content)
                            img = Image.open(image_data)
                            img_width, img_height = img.size
                            aspect_ratio = img_height / img_width
                            width = Inches(5)  # Ajustez la largeur de l'image selon vos besoins
                            height = width * aspect_ratio
                            left = top = Inches(1)  # Ajustez la position de l'image selon vos besoins
                            slide.shapes.add_picture(image_data, left, top, width=width, height=height)
                        else:
                            st.write("Une erreur s'est produite lors du téléchargement de l'image.")

                        # Diapositive avec des points à puces
                        slide_layout = prs.slide_layouts[1]  # Utiliser un layout adapté
                        slide = prs.slides.add_slide(slide_layout)

                        content_shape = slide.placeholders[1]
                        content_shape.text = bullet_points
                        content_shape.text_frame.paragraphs[0].font.size = Pt(12)

                        # Sauvegarder la présentation dans un objet BytesIO
                        ppt_bytes = BytesIO()
                        prs.save(ppt_bytes)
                        ppt_bytes.seek(0)

                        # Télécharger la présentation
                        st.download_button(
                            label="Télécharger PowerPoint",
                            data=ppt_bytes,
                            file_name="presentation.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                        st.write("PowerPoint créé avec succès! 🎉")

                    else:
                        st.write("An error occurred while generating the image.")
                else:
                        st.write("Description generation result is empty.")
                
    cancel_button = st.button('Cancel')
        
    if cancel_button:
        st.stop()
            
if __name__ == "__main__":
    main()

